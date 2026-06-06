#!/usr/bin/env python3
# analyze_more.py — design.md에 3개 섹션 추가: 이미지 내용 분석 / 마스터·플레이스홀더 / 그리드·간격·여백
# 사용: python analyze_more.py <a.pptx> [b.pptx ...]
import sys, os, json, collections, statistics
from pptx import Presentation

def load_elems(pptx):
    j = os.path.splitext(os.path.abspath(pptx))[0] + ".elements.json"
    return json.load(open(j, encoding="utf-8")) if os.path.exists(j) else None

# ---------- 1) 이미지 내용 분석 (형식·배치 기반 추정) ----------
def images_section(d):
    W, H = d["canvas_in"]; imgs = d["images"]
    roles, fmts = collections.Counter(), collections.Counter()
    for im in imgs:
        w, h, ct = im.get("w"), im.get("h"), (im.get("type") or "").lower()
        if "jpeg" in ct or "jpg" in ct: fmts["실사 추정(JPEG)"] += 1
        elif "png" in ct: fmts["PNG(일러스트·아이콘·스크린샷)"] += 1
        elif "emf" in ct or "wmf" in ct: fmts["벡터(EMF/WMF)"] += 1
        elif "svg" in ct: fmts["벡터(SVG)"] += 1
        elif "gif" in ct: fmts["GIF"] += 1
        elif ct: fmts[ct] += 1
        if w and h:
            if w >= 0.72*W and h >= 0.72*H: roles["배경(풀블리드)"] += 1
            elif min(w, h) < 0.1: roles["라인/구분선"] += 1
            elif max(w, h) < 1.0: roles["아이콘·장식"] += 1
            elif max(w, h) >= 6: roles["히어로·대형"] += 1
            else: roles["콘텐츠 이미지"] += 1
    total = len(imgs)
    jpeg = fmts.get("실사 추정(JPEG)", 0)
    ratio = f"{round(jpeg*100/total)}%" if total else "—"
    rl = " · ".join(f"{k} {v}" for k, v in roles.most_common())
    fl = " · ".join(f"{k} {v}" for k, v in fmts.most_common())
    return f"""## 이미지 내용 분석 (형식·배치 기반 추정 · 총 {total}개)
- 역할 분포: {rl or '(없음)'}
- 형식 분포: {fl or '(없음)'}
- **실사(JPEG) 추정 비율: {ratio}** (PNG에 사진·스크린샷이 섞일 수 있어 추정치. 피사체 단위 라벨링은 Vision 분석 필요)
"""

# ---------- 2) 마스터·플레이스홀더 ----------
def master_section(path):
    prs = Presentation(path)
    masters = list(prs.slide_masters)
    lay_lines = []
    for mi, m in enumerate(masters, 1):
        for lay in m.slide_layouts:
            try: nph = len(lay.placeholders)
            except Exception: nph = 0
            lay_lines.append(f"{lay.name}({nph}ph)")
    used = collections.Counter()
    for sl in prs.slides:
        try:
            for ph in sl.placeholders:
                used[str(ph.placeholder_format.type).split()[0]] += 1
        except Exception:
            pass
    used_str = ", ".join(f"{k}×{v}" for k, v in used.most_common()) if used else "**플레이스홀더 미사용** (자유 도형으로 직접 디자인)"
    lay_counter = collections.Counter(lay_lines)
    lay_str = ", ".join(f"{k}×{v}" if v > 1 else k for k, v in lay_counter.most_common()[:12])
    return f"""## 마스터·플레이스홀더
- 슬라이드 마스터: {len(masters)}개
- 정의된 레이아웃({len(lay_lines)}): {lay_str}
- 슬라이드의 실제 플레이스홀더 사용: {used_str}
"""

# ---------- 3) 그리드·간격·여백 ----------
def cluster(vals, tol=0.15):
    vals = sorted(vals); out = []
    for v in vals:
        if out and abs(v - out[-1][0]) <= tol:
            out[-1] = ((out[-1][0]*out[-1][1] + v) / (out[-1][1]+1), out[-1][1]+1)
        else:
            out.append((v, 1))
    return sorted(out, key=lambda c: -c[1])

def grid_section(d):
    W, H = d["canvas_in"]
    lefts, rights, tops, bottoms, gaps = [], [], [], [], []
    for sl in d["slides"]:
        cont = [s for s in sl["shapes"] if s.get("x") is not None and s.get("w")
                and -0.5 <= s["x"] and s["x"]+s["w"] <= W*1.03
                and s["w"] < W*0.72 and (s.get("h") or 0) < H*0.72]
        for s in cont:
            lefts.append(round(s["x"], 2)); rights.append(round(s["x"]+s["w"], 2))
            tops.append(round(s["y"], 2)); bottoms.append(round(s["y"]+s["h"], 2))
        row = sorted(cont, key=lambda s: s["x"])
        for a, b in zip(row, row[1:]):
            g = b["x"] - (a["x"]+a["w"])
            if 0.05 < g < W*0.4: gaps.append(round(g, 2))
    def fmt(cl):
        return ", ".join(f'{round(v,2)}"(×{c})' for v, c in cl[:5]) or "—"
    Lmargin = min((v for v, c in cluster(lefts) if c >= max(3, len(d["slides"])*0.25)), default=None)
    Rextent = max((v for v, c in cluster(rights) if c >= max(3, len(d["slides"])*0.25)), default=None)
    Tmargin = min((v for v, c in cluster(tops) if c >= max(3, len(d["slides"])*0.25)), default=None)
    gap_med = round(statistics.median(gaps), 2) if gaps else None
    gap_mode = collections.Counter(gaps).most_common(1)[0][0] if gaps else None
    return f"""## 그리드·간격·여백 (단위 inch · 캔버스 {W}×{H})
- 좌측 정렬 기준선(빈도순): {fmt(cluster(lefts))}
- 우측 종료선(빈도순): {fmt(cluster(rights))}
- 상단 정렬선(빈도순): {fmt(cluster(tops))}
- 추정 **좌 마진 ≈ {round(Lmargin,2) if Lmargin is not None else '—'}"** · 우 종료 ≈ {round(Rextent,2) if Rextent is not None else '—'}" · 상 마진 ≈ {round(Tmargin,2) if Tmargin is not None else '—'}"
- 요소 간 가로 간격: 중앙값 {gap_med}" · 최빈 {gap_mode}"
"""

MARK = "## 이미지 내용 분석"

def main():
    for p in sys.argv[1:]:
        try:
            d = load_elems(p)
            if d is None:
                print(f"SKIP {p}: elements.json 없음 (먼저 extract_layout.py 실행)", file=sys.stderr); continue
            block = images_section(d) + "\n" + master_section(p) + "\n" + grid_section(d)
            dm = os.path.splitext(os.path.abspath(p))[0] + ".design.md"
            cur = open(dm, encoding="utf-8").read() if os.path.exists(dm) else ""
            if MARK in cur:
                cur = cur[:cur.index(MARK)].rstrip() + "\n\n"
            open(dm, "w", encoding="utf-8").write(cur.rstrip() + "\n\n" + block)
            print(f"OK  {os.path.basename(dm)}")
        except Exception as e:
            print(f"FAIL {p}: {e}", file=sys.stderr)

if __name__ == "__main__":
    main()
