#!/usr/bin/env python3
# extract_layout.py — pptx의 템플릿·이미지·내용·요소 위치/배치를 추출.
#   - <파일명>.elements.json : 슬라이드별 전 도형의 종류·좌표·크기·텍스트·이미지(정밀)
#   - 콘솔/리턴: design.md에 붙일 "슬라이드 구조" 요약 마크다운
# 사용: python extract_layout.py <a.pptx> [b.pptx ...]
import sys, os, json, collections
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def emu_in(v):
    return round(v / 914400, 2) if v is not None else None

def kind(sh):
    try:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE: return "image"
        if sh.has_table: return "table"
        if sh.has_chart: return "chart"
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP: return "group"
    except Exception:
        pass
    if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip(): return "text"
    return "shape"

def zone(cx, cy, W, H):
    col = "L" if cx < W/3 else ("C" if cx < 2*W/3 else "R")
    row = "상" if cy < H/3 else ("중" if cy < 2*H/3 else "하")
    return col + row

def analyze(path):
    prs = Presentation(path)
    W, H = emu_in(prs.slide_width), emu_in(prs.slide_height)
    Wf, Hf = prs.slide_width, prs.slide_height
    layouts = collections.Counter()
    kinds = collections.Counter()
    images = []
    slides = []
    for idx, sl in enumerate(prs.slides, 1):
        lname = ""
        try: lname = sl.slide_layout.name
        except Exception: pass
        layouts[lname] += 1
        shapes = []
        for sh in sl.shapes:
            k = kind(sh); kinds[k] += 1
            x, y = emu_in(sh.left), emu_in(sh.top)
            w, h = emu_in(sh.width), emu_in(sh.height)
            z = ""
            if None not in (sh.left, sh.top, sh.width, sh.height):
                z = zone(sh.left + sh.width/2, sh.top + sh.height/2, Wf, Hf)
            rec = {"type": k, "zone": z, "x": x, "y": y, "w": w, "h": h}
            if k == "text":
                t = sh.text_frame.text.strip().replace("\n", " ")
                rec["text"] = t[:60]
            if k == "image":
                try:
                    img = sh.image
                    fn = os.path.basename(img.filename or f"img.{img.ext}")
                    rec["image"] = fn
                    images.append({"slide": idx, "file": fn, "type": img.content_type, "x": x, "y": y, "w": w, "h": h, "zone": z})
                except Exception:
                    rec["image"] = "?"
                    images.append({"slide": idx, "file": "?", "x": x, "y": y, "w": w, "h": h, "zone": z})
            shapes.append(rec)
        slides.append({"slide": idx, "layout": lname, "shapes": shapes})
    return {"canvas_in": [W, H], "n_slides": len(slides), "layouts": dict(layouts),
            "kind_counts": dict(kinds), "n_images": len(images), "images": images, "slides": slides}

def one_line(sl):
    parts = []
    order = {"text": 0, "image": 1, "table": 2, "chart": 3, "shape": 4, "group": 5}
    for sh in sorted(sl["shapes"], key=lambda s: (s.get("zone", "z"), order.get(s["type"], 9))):
        if sh["type"] == "image":
            parts.append(f"이미지({sh.get('zone','')} {sh.get('w','?')}×{sh.get('h','?')})")
        elif sh["type"] == "text" and sh.get("text"):
            parts.append(f"텍스트[{sh['zone']}]:“{sh['text'][:14]}”")
        elif sh["type"] in ("table", "chart"):
            parts.append(f"{ {'table':'표','chart':'차트'}[sh['type']] }({sh.get('zone','')})")
    return " · ".join(parts[:6]) if parts else "(빈 슬라이드/배경만)"

def summary_md(d):
    lay = ", ".join(f"{k or '(이름없음)'}×{v}" for k, v in sorted(d["layouts"].items(), key=lambda kv: -kv[1]))
    kc = d["kind_counts"]
    kinds = " · ".join(f"{ko} {kc[k]}" for k, ko in [("text","텍스트"),("image","이미지"),("shape","도형"),("table","표"),("chart","차트"),("group","그룹")] if kc.get(k))
    lines = "\n".join(f"  - S{sl['slide']}: {one_line(sl)}" for sl in d["slides"][:16])
    more = "" if d["n_slides"] <= 16 else f"\n  - … 외 {d['n_slides']-16}장"
    return f"""## 슬라이드 구조 (템플릿·이미지·내용·요소배치)
- 사용 레이아웃(템플릿): {lay}
- 요소 타입 분포: {kinds}
- 이미지: 총 **{d['n_images']}개** (정밀 위치·크기는 `<파일명>.elements.json`)
- 위치 표기: L/C/R(좌·중·우) × 상/중/하. 단위 inch.
- 슬라이드별 요소배치(요약):
{lines}{more}
"""

def main():
    if len(sys.argv) < 2:
        print("usage: python extract_layout.py <pptx ...>", file=sys.stderr); return 1
    for p in sys.argv[1:]:
        try:
            d = analyze(p)
            base = os.path.splitext(os.path.abspath(p))[0]
            with open(base + ".elements.json", "w", encoding="utf-8") as f:
                json.dump(d, f, ensure_ascii=False, indent=1)
            # design.md 에 구조 섹션 append (중복 방지: 마커로 교체)
            dm = base + ".design.md"
            marker = "## 슬라이드 구조"
            block = summary_md(d)
            if os.path.exists(dm):
                cur = open(dm, encoding="utf-8").read()
                if marker in cur:
                    cur = cur[:cur.index(marker)].rstrip() + "\n\n"
                open(dm, "w", encoding="utf-8").write(cur.rstrip() + "\n\n" + block)
            else:
                open(dm, "w", encoding="utf-8").write(block)
            print(f"OK  {os.path.basename(base)} (slides={d['n_slides']}, images={d['n_images']})")
        except Exception as e:
            print(f"FAIL {p}: {e}", file=sys.stderr)
    return 0

if __name__ == "__main__":
    sys.exit(main())
